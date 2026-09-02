# -*- coding: utf-8 -*-
"""ppt_maker.md 기준으로 3_Credit_Card_Fraud_Detection.pptx를 수정해
3_Credit_Card_Fraud_Detection(1).pptx로 저장한다."""
import os
import tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

SRC = "report/4차/3_Credit_Card_Fraud_Detection.pptx"
OUT = "report/4차/3_Credit_Card_Fraud_Detection(1).pptx"
ASSETS = "report/4차/report_assets/"
TMP = tempfile.gettempdir()  # 그림 2-5 왼쪽 패널만 잘라낸 중간 파일 위치

BLUE = "3970CF"
BLUE2 = "2664B4"
DARK = "222222"
LIGHT = "F4F4F4"
WHITE = "FFFFFF"
RED = "C00000"

FONT = "JASO Sans"
FONT_B = "JASO Sans Bold"

CONTENT_LEFT = 1028700
CONTENT_W = 16230600


def rgb(hexstr):
    return RGBColor.from_string(hexstr)


# ---------------------------------------------------------------- helpers (build_ppt.py 동일)
def _set_alpha(color_format, alpha_pct):
    srgb = color_format._xFill.find(qn('a:srgbClr'))
    a = etree.SubElement(srgb, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))


def add_textbox(slide, text, left, top, width, height, size=18, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, font=None, alpha=None, anchor=None, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font or (FONT_B if bold else FONT)
        r.font.color.rgb = rgb(color)
        if alpha is not None:
            _set_alpha(r.font.color, alpha)
    return box


def add_rect(slide, left, top, width, height, fill=LIGHT, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
    shp.shadow.inherit = False
    return shp


def add_stat_card(slide, number, label, left, top, width, height):
    # 배경(LIGHT)과 같은 색이면 블럭이 보이지 않으므로 흰 면 + 옅은 테두리로 처리
    add_rect(slide, left, top, width, height, fill=WHITE, line="D9D9D9")
    add_textbox(slide, number, left, top + 300000, width, 620000, size=32, bold=True,
                color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left + 60000, top + height - 620000, width - 120000, 380000,
                size=13, color=DARK, align=PP_ALIGN.CENTER)


def add_picture_fit(slide, path, left, top, max_w, max_h):
    iw, ih = Image.open(path).size
    ratio = iw / ih
    w, h = max_w, int(max_w / ratio)
    if h > max_h:
        h, w = max_h, int(max_h * ratio)
    slide.shapes.add_picture(path, left + (max_w - w) // 2, top + (max_h - h) // 2,
                             width=w, height=h)


def add_table_styled(slide, data, left, top, width, height, col_widths=None,
                     header_size=13, body_size=12, highlight_rows=None, align_center_cols=None):
    rows, cols = len(data), len(data[0])
    highlight_rows = highlight_rows or set()
    align_center_cols = align_center_cols if align_center_cols is not None else set(range(1, cols))
    gshape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gshape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        table.rows[r].height = int(height / rows)
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Emu(60000)
            cell.margin_top = cell.margin_bottom = Emu(20000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (c in align_center_cols or r == 0) else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.size = Pt(header_size if r == 0 else body_size)
            run.font.name = FONT_B if (r == 0 or r in highlight_rows) else FONT
            run.font.bold = (r == 0) or (r in highlight_rows)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = rgb(BLUE)
                run.font.color.rgb = rgb(WHITE)
            elif r in highlight_rows:
                cell.fill.fore_color.rgb = rgb("D6E4FA")
                run.font.color.rgb = rgb(BLUE2)
            else:
                cell.fill.fore_color.rgb = rgb(WHITE if r % 2 else LIGHT)
                run.font.color.rgb = rgb(DARK)
    return gshape


def set_bg(slide, hexcolor):
    cSld = slide._element.find(qn('p:cSld'))
    old = cSld.find(qn('p:bg'))
    if old is not None:
        cSld.remove(old)
    cSld.insert(0, etree.fromstring(
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="%s"/></a:solidFill></p:bgPr></p:bg>' % hexcolor))


# ---------------------------------------------------------------- 편집용 helper
def clear_slide(slide):
    for sp in list(slide.shapes):
        sp._element.getparent().remove(sp._element)


def find_by_prefix(slide, prefix):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith(prefix):
            return sh
    raise KeyError("prefix %r not found" % prefix)


def retitle(slide, prefix, new_text):
    """제목 텍스트박스 내용만 교체(서식 유지)."""
    p = find_by_prefix(slide, prefix).text_frame.paragraphs[0]
    runs = p.runs
    for r in runs[1:]:
        p._p.remove(r._r)
    runs[0].text = new_text


def sub_text(slide, mapping):
    """슬라이드 전체 run에서 부분 문자열 치환."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                for a, b in mapping.items():
                    if a in r.text:
                        r.text = r.text.replace(a, b)


def outline_row(table, row_idx, color=RED, width_emu=28575):
    """행 바깥 테두리를 굵은 빨간 선으로 처리."""
    ncols = len(table.columns)
    for c in range(ncols):
        tc = table.cell(row_idx, c)._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))
        edges = ['a:lnT', 'a:lnB']
        if c == 0:
            edges.append('a:lnL')
        if c == ncols - 1:
            edges.append('a:lnR')
        # tcPr 자식 순서는 lnL, lnR, lnT, lnB 이므로 역순으로 맨 앞에 삽입한다
        for tag in ['a:lnB', 'a:lnT', 'a:lnR', 'a:lnL']:
            if tag not in edges:
                continue
            for old in tcPr.findall(qn(tag)):
                tcPr.remove(old)
            ln = etree.fromstring(
                '<%s xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                'w="%d" cap="flat" cmpd="sng" algn="ctr">'
                '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
                '<a:prstDash val="solid"/></%s>' % (tag, width_emu, color, tag))
            tcPr.insert(0, ln)


def reorder(prs, order):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for e in ids:
        lst.remove(e)
    for i in order:
        lst.append(ids[i])


def crop_left(src, dst, right_px):
    im = Image.open(src)
    im.crop((0, 0, right_px, im.size[1])).save(dst)
    return dst


# =====================================================================
prs = Presentation(SRC)
S = prs.slides
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- S3  2.1 프로젝트 개요
s3 = S[2]
clear_slide(s3)
set_bg(s3, LIGHT)
add_textbox(s3, "2.1 프로젝트 개요", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(
    s3,
    "28개의 PCA(주성분분석) 변환 변수(V1~V28)와 거래 시각(Time), 거래 금액(Amount)으로\n"
    "사기 거래 여부(Class)를 예측하는 이진 분류 문제\n"
    "샘플링 기법 3가지(언더샘플링 · SMOTE · SMOTE-ENN)와 모델 5가지"
    "(RandomForest · LightGBM · XGBoost · GradientBoost · LogisticRegression)의 성능을 비교한다",
    CONTENT_LEFT + 800000, 1600000, CONTENT_W - 1600000, 1600000,
    size=17, color=DARK, alpha=80, align=PP_ALIGN.CENTER, line_spacing=1.3)
stats = [
    ("284,807", "전체 거래 건수"),
    ("31", "컬럼 수 (Time+V1~V28+Amount+Class)"),
    ("0.173%", "사기(Class=1) 비율"),
    ("0건", "결측치"),
]
card_w, card_h, card_top = 3600000, 1900000, 3800000
gap = (CONTENT_W - 4 * card_w) // 3
for i, (num, label) in enumerate(stats):
    add_stat_card(s3, num, label, CONTENT_LEFT + i * (card_w + gap), card_top, card_w, card_h)

# ---------------------------------------------------------------- S4  2.2 EDA - Amount
amount_left = crop_left(ASSETS + "credit_amount.png",
                        os.path.join(TMP, "credit_amount_left.png"), 755)
s4 = S[3]
clear_slide(s4)
set_bg(s4, LIGHT)
add_textbox(s4, "2.2 EDA — 거래 금액 분포 (Amount)", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_picture_fit(s4, amount_left, CONTENT_LEFT, 1700000, 9200000, 6900000)
add_textbox(s4, "그림 2-5(좌). Amount 원본 분포 — 왜도(skew) 16.98",
            CONTENT_LEFT, 8800000, 9200000, 500000,
            size=14, color=DARK, alpha=85, align=PP_ALIGN.CENTER)
bullets4 = [
    "평균 $88.35 > 중앙값 $22.00 — 소액 거래가 대부분이고 극소수 고액 거래가 평균을 끌어올린다",
    "왜도(skewness) 16.98 — 오른쪽으로 극단적으로 긴 꼬리(right-skewed) 분포",
    "사기 거래: 평균 $122.21(+38%) vs 중앙값 $9.25(−58%) — 소액 \"테스트 결제\"와 일부 고액 사기가 공존",
    "→ log1p 변환으로 왜도 16.98 → 0.16, Amount_Scaled 컬럼 생성 근거 (§2.4)",
]
for i, b in enumerate(bullets4):
    add_textbox(s4, "•  " + b, 10800000, 2400000 + i * 1150000,
                CONTENT_W - (10800000 - CONTENT_LEFT), 1050000,
                size=16, color=DARK, alpha=85, line_spacing=1.2)

# ---------------------------------------------------------------- S5  2.2 EDA - Time
s5 = S[4]
clear_slide(s5)
set_bg(s5, LIGHT)
add_textbox(s5, "2.2 EDA — 클래스와 시간대 패턴 (Time)", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_picture_fit(s5, ASSETS + "credit_time_pattern.png", CONTENT_LEFT, 1600000, CONTENT_W, 4700000)
add_textbox(s5,
            "그림 2-4. 시간대별 거래 건수 — 정상 거래(왼쪽)는 새벽 0~6시에 뚜렷이 감소하지만, "
            "사기 거래(오른쪽)는 이 시간대에도 상대적으로 활발하다.",
            CONTENT_LEFT, 6500000, CONTENT_W, 600000,
            size=14, color=DARK, alpha=85, align=PP_ALIGN.CENTER)
bullets5 = [
    "Time은 첫 거래 이후 경과 초 — (Time/3600) % 24로 환산해 하루 중 거래 시간대를 확인",
    "새벽 0~6시 사기 비중 25.20% vs 정상 8.37% — 약 3배, 소유자가 잠든 시간의 도난·탈취 거래 패턴과 부합",
    "다만 Time은 실제 시계(time-of-day)가 아닌 수집 시작 후 경과 초 → 피처로 남기지 않고 제거 (§2.4)",
]
for i, b in enumerate(bullets5):
    add_textbox(s5, "•  " + b, CONTENT_LEFT, 7350000 + i * 720000, CONTENT_W, 680000,
                size=16, color=DARK, alpha=85, line_spacing=1.2)

# ---------------------------------------------------------------- S6  2.2 EDA - Class (신규)
s6 = S.add_slide(BLANK)
set_bg(s6, LIGHT)
add_textbox(s6, "2.2 EDA — 클래스 분포 (Class)", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
half_w = (CONTENT_W - 400000) // 2
right_x = CONTENT_LEFT + half_w + 400000
add_picture_fit(s6, ASSETS + "credit_target.png", CONTENT_LEFT, 1420000, half_w, 5300000)
add_picture_fit(s6, ASSETS + "credit_corr.png", right_x, 1420000, half_w, 5300000)
add_textbox(s6, "그림 2-2. Class 분포(로그 스케일) — 정상 284,315건(99.827%) vs 사기 492건(0.173%), 약 578:1",
            CONTENT_LEFT, 6800000, half_w, 800000,
            size=14, color=DARK, alpha=85, align=PP_ALIGN.CENTER, line_spacing=1.2)
add_textbox(s6, "그림 2-3. 31개 변수 전체의 피어슨 상관계수 히트맵 — V1~V28은 PCA 축이라 서로 직교, "
                "색이 뚜렷이 갈리는 곳은 Time·Amount·Class 행뿐",
            right_x, 6800000, half_w, 800000,
            size=14, color=DARK, alpha=85, align=PP_ALIGN.CENTER, line_spacing=1.2)
add_textbox(s6, "Santander(24:1)보다 100배 이상 심한 불균형이다. 이 정도 비율에서는 ROC-AUC조차 "
                "낙관적으로 보일 수 있어, ROC-AUC와 AUPRC를 함께 본다.",
            CONTENT_LEFT, 7750000, half_w, 1000000,
            size=15, color=DARK, alpha=90, line_spacing=1.3)
add_textbox(s6, "Class와 V14와의 상관관계(−0.303)가 높은 것으로 추정되어 V14의 이상치 제거 계획.",
            right_x, 7750000, half_w, 1000000,
            size=15, color=DARK, alpha=90, line_spacing=1.3)
add_textbox(s6, "극단적 불균형(578:1) — ROC-AUC만으로 부족하므로 AUPRC를 함께 확인 (§2.3)",
            CONTENT_LEFT, 9000000, CONTENT_W, 600000,
            size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- S7  2.3 결과 지표 정의
s7 = S[5]
retitle(s7, "3.2 결과 지표", "2.3 결과 지표 정의")
tbl7 = next(sh for sh in s7.shapes if sh.has_table).table
outline_row(tbl7, 1)   # Recall
outline_row(tbl7, 4)   # ROC-AUC

# ---------------------------------------------------------------- S8  2.3 재현율 강조 이유
retitle(S[6], "3.2 왜 재현율", "2.3 왜 재현율(Recall)을 우선하는가")

# ---------------------------------------------------------------- S9  2.4 데이터 전처리
retitle(S[7], "3.3 데이터 전처리", "2.4 데이터 전처리 & 클리닝")

# ---------------------------------------------------------------- S10 2.5 샘플링 비교
s10 = S[8]
clear_slide(s10)
set_bg(s10, LIGHT)
add_textbox(s10, "2.5 왜 리샘플링이 필요한가 — 샘플링 전략 3종 비교", CONTENT_LEFT, 620000, CONTENT_W, 700000,
            size=26, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tbl_a = [
    ["기법", "방식", "학습 데이터 변화"],
    ["언더샘플링", "정상 거래를 사기 건수만큼 무작위로 축소", "199,362건 → 684건 (342/342)"],
    ["SMOTE (채택)", "사기 샘플 주변에 합성 데이터 생성해 오버샘플링", "199,362건 → 398,040건 (199,020/199,020)"],
    ["SMOTE-ENN", "SMOTE 이후 경계가 모호한 샘플을 ENN으로 추가 제거", "199,362건 → 397,717건 (199,020/198,697)"],
]
add_table_styled(s10, tbl_a, CONTENT_LEFT, 1500000, CONTENT_W, 1900000,
                 col_widths=[Emu(2400000), Emu(7200000), Emu(6630600)],
                 header_size=14, body_size=13, align_center_cols={0}, highlight_rows={2})
tbl_b = [
    ["기법", "오차 행렬 (TN · FP / FN · TP)", "정확도", "정밀도", "재현율", "F1", "AUC", "AUPRC"],
    ["언더샘플링", "82,703 · 2,592 / 17 · 130", "0.9695", "0.0478", "0.8844", "0.0906", "0.9777", "0.6689"],
    ["오버샘플링 (SMOTE, 채택)", "85,278 · 17 / 30 · 117", "0.9994", "0.8731", "0.7959", "0.8327", "0.9697", "0.8274"],
    ["오버샘플링 (SMOTE-ENN)", "85,254 · 41 / 27 · 120", "0.9992", "0.7453", "0.8163", "0.7792", "0.9739", "0.8059"],
]
add_table_styled(s10, tbl_b, CONTENT_LEFT, 3600000, CONTENT_W, 1900000,
                 col_widths=[Emu(3400000), Emu(4230600), Emu(1440000), Emu(1440000),
                             Emu(1440000), Emu(1280000), Emu(1440000), Emu(1560000)],
                 header_size=13, body_size=12, align_center_cols={0}, highlight_rows={2})
add_picture_fit(s10, ASSETS + "credit_smote_scatter.png", CONTENT_LEFT, 5600000, CONTENT_W, 3100000)
add_textbox(s10, "SMOTE는 정상 거래 정보를 버리지 않으면서 사기 샘플을 원본의 약 582배로 증강",
            CONTENT_LEFT, 8800000, CONTENT_W, 500000,
            size=15, color=DARK, alpha=85, align=PP_ALIGN.CENTER)
add_textbox(s10, "채택: SMOTE — 언더샘플링은 정밀도 0.0478로 오탐 과다, SMOTE-ENN은 SMOTE 대비 이점 미미 (§2.8)",
            CONTENT_LEFT, 9350000, CONTENT_W, 500000,
            size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- S11 2.6 GridSearchCV
s11 = S[9]
retitle(s11, "3.5 GridSearchCV", "2.6 GridSearchCV 튜닝 — RandomForest")
sub_text(s11, {"(§3.6.2)": "(§2.7)"})

# ---------------------------------------------------------------- S12/S13 성능 지표 · 파이프라인
retitle(S[11], "3.6 성능 지표 결과", "2.7 성능 지표 결과 — RandomForest + SMOTE")
retitle(S[10], "3.6 파이프라인 요약", "2.7 파이프라인 요약")

# ---------------------------------------------------------------- S14 2.8 모델별 비교
retitle(S[12], "3.6 모델별 비교", "2.8 모델별 비교 (SMOTE 적용)")

# ---------------------------------------------------------------- S15 종합 결론 / S16 감사합니다
retitle(S[13], "3장 종합 결론", "2장 종합 결론")
sub_text(S[14], {"3장 Credit Card Fraud Detection": "2장 Credit Card Fraud Detection"})

# ---------------------------------------------------------------- 순서 재배치
# 현재: 0표지 1목차 2개요 3Amount 4Time 5지표 6재현율 7전처리 8샘플링 9Grid
#       10파이프라인 11성능 12모델별 13결론 14감사 15Class(신규)
reorder(prs, [0, 1, 2, 3, 4, 15, 5, 6, 7, 8, 9, 11, 10, 12, 13, 14])

prs.save(OUT)
print("SAVED: %s  (slides = %d)" % (OUT, len(prs.slides._sldIdLst)))
