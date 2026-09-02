# -*- coding: utf-8 -*-
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from PIL import Image

SRC_TEMPLATE = "ppt_예시.pptx"
OUT_PATH = "report/4차/3_Credit_Card_Fraud_Detection.pptx"
IMG_DIR = "scratch_section3_images/"

BLUE = "3970CF"
BLUE2 = "2664B4"
DARK = "222222"
LIGHT = "F4F4F4"
LIGHTER = "FCFCFC"
WHITE = "FFFFFF"

SLIDE_W = 18288000
SLIDE_H = 10287000
FONT = "JASO Sans"
FONT_B = "JASO Sans Bold"


def rgb(hexstr):
    return RGBColor.from_string(hexstr)


# ---------------------------------------------------------------- slide ops
def duplicate_slide(prs, index):
    """Deep-copy slide at `index`, append at end, remap relationships. Returns new slide."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    # remove any placeholder shapes python-pptx auto-added for the layout
    for sp in list(dest.shapes):
        sp._element.getparent().remove(sp._element)

    # copy background
    src_cSld = source._element.find(qn('p:cSld'))
    dest_cSld = dest._element.find(qn('p:cSld'))
    old_bg = dest_cSld.find(qn('p:bg'))
    if old_bg is not None:
        dest_cSld.remove(old_bg)
    src_bg = src_cSld.find(qn('p:bg'))
    if src_bg is not None:
        dest_cSld.insert(0, copy.deepcopy(src_bg))

    # relationship remap
    rid_map = {}
    for rId, rel in source.part.rels.items():
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue
        if rel.is_external:
            new_rid = dest.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rId] = new_rid

    for shape in source.shapes:
        new_el = copy.deepcopy(shape._element)
        if rid_map:
            xml_str = etree.tostring(new_el).decode('utf-8')
            for old, new in sorted(rid_map.items(), key=lambda kv: -len(kv[0])):
                xml_str = xml_str.replace(f'"{old}"', f'"{new}"')
            new_el = etree.fromstring(xml_str.encode('utf-8'))
        dest.shapes._spTree.append(new_el)

    return dest


def remove_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"shape {name!r} not found")


def clean_set_text(shape, text):
    """Replace a shape's entire text with `text`, collapsing any extra runs/paragraphs
    down to the first run's formatting (safe for original multi-run/multi-paragraph shapes)."""
    txBody = shape.text_frame._txBody
    ps = txBody.findall(qn('a:p'))
    first_p = ps[0]
    for p in ps[1:]:
        txBody.remove(p)
    runs = first_p.findall(qn('a:r'))
    if runs:
        for r in runs[1:]:
            first_p.remove(r)
        runs[0].find(qn('a:t')).text = text
    else:
        r_xml = (
            '<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:t>{text}</a:t></a:r>'
        )
        first_p.append(etree.fromstring(r_xml))


def set_para_text_collapse(shape, text, para_idx=0):
    """Like clean_set_text but targets one paragraph among several, leaving other
    paragraphs untouched (used for the 2-line cover title)."""
    txBody = shape.text_frame._txBody
    p = txBody.findall(qn('a:p'))[para_idx]
    runs = p.findall(qn('a:r'))
    for r in runs[1:]:
        p.remove(r)
    if runs:
        runs[0].find(qn('a:t')).text = text
    return p.find(qn('a:r'))


# ---------------------------------------------------------------- component builders
def add_textbox(slide, text, left, top, width, height, size=18, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, font=None, alpha=None, anchor=None, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font or (FONT_B if bold else FONT)
        r.font.color.rgb = rgb(color)
        if alpha is not None:
            _set_alpha(r.font.color, alpha)
    return box


def _set_alpha(color_format, alpha_pct):
    """alpha_pct: 0-100"""
    srgb = color_format._xFill.find(qn('a:srgbClr'))
    a = etree.SubElement(srgb, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))


def add_title(slide, text, top=520000, size=30, color=BLUE):
    return add_textbox(slide, text, 914400, top, SLIDE_W - 2 * 914400, 700000,
                        size=size, bold=True, color=color, align=PP_ALIGN.LEFT, font=FONT_B)


def add_rect(slide, left, top, width, height, fill=LIGHT, line=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
    shp.shadow.inherit = False
    return shp


def add_stat_card(slide, number, label, left, top, width, height):
    add_rect(slide, left, top, width, height, fill=LIGHT)
    add_textbox(slide, number, left, top + 130000, width, 620000, size=32, bold=True,
                color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left + 60000, top + height - 400000, width - 120000, 380000,
                size=13, bold=False, color=DARK, align=PP_ALIGN.CENTER)
    return


def add_picture_fit(slide, path, left, top, max_w, max_h):
    im = Image.open(path)
    iw, ih = im.size
    ratio = iw / ih
    w = max_w
    h = int(w / ratio)
    if h > max_h:
        h = max_h
        w = int(h * ratio)
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    slide.shapes.add_picture(path, x, y, width=w, height=h)
    return x, y, w, h


def add_table_styled(slide, data, left, top, width, height, col_widths=None,
                      header_size=13, body_size=12, highlight_rows=None, align_center_cols=None):
    """data: list of rows (list of str), first row = header."""
    rows = len(data)
    cols = len(data[0])
    highlight_rows = highlight_rows or set()
    align_center_cols = align_center_cols if align_center_cols is not None else set(range(1, cols))
    gshape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gshape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        row_h = int(height / rows)
        table.rows[r].height = row_h
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Emu(60000)
            cell.margin_right = Emu(60000)
            cell.margin_top = Emu(20000)
            cell.margin_bottom = Emu(20000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (c in align_center_cols or r == 0) else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.size = Pt(header_size if r == 0 else body_size)
            run.font.name = FONT_B if (r == 0 or r in highlight_rows) else FONT
            run.font.bold = (r == 0) or (r in highlight_rows)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = rgb(BLUE)
                run.font.color.rgb = rgb(WHITE)
            elif r in highlight_rows:
                cell.fill.fore_color.rgb = rgb("D6E4FA")
                run.font.color.rgb = rgb(BLUE2)
            else:
                cell.fill.fore_color.rgb = rgb(WHITE if r % 2 else LIGHT)
                run.font.color.rgb = rgb(DARK)
    return gshape


def set_bg(slide, hexcolor):
    cSld = slide._element.find(qn('p:cSld'))
    old = cSld.find(qn('p:bg'))
    if old is not None:
        cSld.remove(old)
    bg_xml = (
        f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:bgPr><a:solidFill><a:srgbClr val="{hexcolor}"/></a:solidFill></p:bgPr></p:bg>'
    )
    cSld.insert(0, etree.fromstring(bg_xml))


def restyle_title(slide, shape_name, text, size_pt=28, color=DARK):
    shape = find_shape(slide, shape_name)
    shape.left = Emu(1028700)
    shape.width = Emu(16230600)
    shape.top = Emu(700000)
    shape.height = Emu(750000)
    p = shape.text_frame.paragraphs[0]
    r = p.runs[0]
    r.text = text
    r.font.size = Pt(size_pt)
    r.font.color.rgb = rgb(color)
    return shape


CONTENT_LEFT = 1028700
CONTENT_W = 16230600
BODY_TOP = 1750000


print("helpers loaded OK")

# =====================================================================
# BUILD DECK
# =====================================================================

prs = Presentation(SRC_TEMPLATE)

# ---------------------------------------------------------------- S1 표지
s1 = duplicate_slide(prs, 0)
clean_set_text(find_shape(s1, "TextBox 2"), "신용카드 사기 거래 탐지 · 이진 분류 · 초극단적 클래스 불균형")
clean_set_text(find_shape(s1, "TextBox 3"), "담당  |  염윤호")
clean_set_text(find_shape(s1, "TextBox 4"), "작성일  |  2026. 9. 2")
title5 = find_shape(s1, "TextBox 5")
set_para_text_collapse(title5, "Credit Card", para_idx=0)
set_para_text_collapse(title5, "Fraud Detection", para_idx=1)
sub6 = find_shape(s1, "TextBox 6")
clean_set_text(sub6, "MACHINE LEARNING CASE STUDY")
sub6.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

# ---------------------------------------------------------------- S2 목차
s2 = duplicate_slide(prs, 1)
agenda_map = {
    "TextBox 2": "프로젝트 개요",
    "TextBox 4": "EDA — 탐색적 데이터 분석",
    "TextBox 10": "결과 지표 정의",
    "TextBox 14": "데이터 전처리",
    "TextBox 3": "모델링 & 샘플링 전략",
    "TextBox 5": "파이프라인 & 성능 검증",
    "TextBox 11": "모델별 비교",
    "TextBox 15": "종합 결론",
}
for name, text in agenda_map.items():
    clean_set_text(find_shape(s2, name), text)

# ---------------------------------------------------------------- S3 프로젝트 개요
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, LIGHT)
add_textbox(s3, "3.1 프로젝트 개요", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(
    s3,
    "28개의 PCA 변환 변수(V1~V28)와 거래 시각(Time), 거래 금액(Amount)으로 사기 거래 여부(Class)를 예측하는 "
    "이진 분류 문제. 담당 모델은 RandomForest + SMOTE이며, 임곗값을 재현율 중심으로 조정한 것이 최종 채택안이다.",
    CONTENT_LEFT + 800000, 1650000, CONTENT_W - 1600000, 1200000,
    size=17, color=DARK, alpha=80, align=PP_ALIGN.CENTER, line_spacing=1.3,
)
stats = [
    ("284,807", "전체 거래 건수"),
    ("31", "컬럼 수 (Time+V1~V28+Amount+Class)"),
    ("0.173%", "사기(Class=1) 비율"),
    ("0건", "결측치"),
]
card_w = 3600000
gap = (CONTENT_W - 4 * card_w) // 3
card_top = 3400000
card_h = 2200000
for i, (num, label) in enumerate(stats):
    left = CONTENT_LEFT + i * (card_w + gap)
    add_stat_card(s3, num, label, left, card_top, card_w, card_h)

print("S1-S3 built OK")

# ---------------------------------------------------------------- S4 EDA - Amount 분포
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, LIGHT)
add_textbox(s4, "3.1 EDA — 거래 금액 분포 (Amount)", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_picture_fit(s4, IMG_DIR + "p17_0.png", CONTENT_LEFT, 1650000, 9600000, 7800000)
bullets4 = [
    "평균 $88.35 > 중앙값 $22.00 — 오른쪽 꼬리 분포 (왜도 16.98)",
    "사기 거래: 평균 $122.21(+38%) vs 중앙값 $9.25(-58%)",
    "소액 \"테스트 결제\"와 일부 고액 사기가 공존하는 양극화 구조",
    "→ Amount에 log1p 변환 적용 근거 (§3.3)",
]
bx4 = 11100000
by4 = 2200000
for i, b in enumerate(bullets4):
    add_textbox(s4, "•  " + b, bx4, by4 + i * 900000, CONTENT_W - (bx4 - CONTENT_LEFT), 820000,
                size=16, color=DARK, alpha=85, line_spacing=1.2)

# ---------------------------------------------------------------- S5 EDA - 클래스 불균형 & 시간대
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, LIGHT)
add_textbox(s5, "3.1 EDA — 클래스 불균형과 시간대 패턴", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
half_w = (CONTENT_W - 400000) // 2
add_picture_fit(s5, IMG_DIR + "p18_0.png", CONTENT_LEFT, 1650000, half_w, 5600000)
add_textbox(s5, "정상 284,315건(99.827%) vs 사기 492건(0.173%) — 약 578:1", CONTENT_LEFT, 7350000,
            half_w, 700000, size=15, color=DARK, alpha=85, align=PP_ALIGN.CENTER, line_spacing=1.2)
add_picture_fit(s5, IMG_DIR + "p20_0.png", CONTENT_LEFT + half_w + 400000, 1650000, half_w, 5600000)
add_textbox(s5, "새벽 0~6시 사기 비중 25.20% vs 정상 8.37% — 약 3배", CONTENT_LEFT + half_w + 400000, 7350000,
            half_w, 700000, size=15, color=DARK, alpha=85, align=PP_ALIGN.CENTER, line_spacing=1.2)
add_textbox(s5, "극단적 불균형에서는 ROC-AUC만으로 부족 → AUPRC를 함께 확인 (§3.2)",
            CONTENT_LEFT, 8250000, CONTENT_W, 500000, size=16, bold=True, color=BLUE,
            align=PP_ALIGN.CENTER)

print("S4-S5 built OK")

# ---------------------------------------------------------------- S6 결과 지표 정의
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, LIGHT)
add_textbox(s6, "3.2 결과 지표 정의", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tbl6 = [
    ["지표", "정의", "이 데이터에서의 역할"],
    ["Recall (재현율)", "실제 사기 중 모델이 찾아낸 비율", "최우선 지표 — 놓친 사기는 금전 손실"],
    ["Precision (정밀도)", "사기로 예측한 것 중 실제 사기 비율", "너무 낮으면 정상 거래 차단 급증"],
    ["F1-Score", "정밀도·재현율의 조화평균", "재현율 편중을 견제하는 균형 지표"],
    ["ROC-AUC", "전체 임곗값 TPR-FPR 트레이드오프", "모델 판별력 비교 (단독 과대평가 위험)"],
    ["AUPRC", "Precision-Recall 곡선의 면적", "극단적 불균형에서 ROC-AUC를 보완하는 핵심 지표"],
]
add_table_styled(s6, tbl6, CONTENT_LEFT, 1900000, CONTENT_W, 4600000,
                  col_widths=[Emu(3200000), Emu(6200000), Emu(6830600)],
                  header_size=16, body_size=15, align_center_cols={0})
add_textbox(s6, "이 프로젝트의 핵심 의도 — 정밀도와 재현율 중 재현율을 우선한다",
            CONTENT_LEFT, 6800000, CONTENT_W, 600000, size=18, bold=True, color=BLUE,
            align=PP_ALIGN.CENTER)

print("S6 built OK")

# ---------------------------------------------------------------- S7 왜 재현율을 우선하는가 (SWOT 2x2 재사용)
s7 = duplicate_slide(prs, 7)
clean_set_text(find_shape(s7, "TextBox 2"), "3.2 왜 재현율(Recall)을 우선하는가")
clean_set_text(find_shape(s7, "TextBox 6"), (
    "카드 사기를 놓치면(FN) 그대로 금전 손실이지만, 정상 거래를 오탐지(FP)해도 "
    "고객 확인 연락 정도의 비용에 그친다 — 그래서 FN을 줄이는 방향으로 임곗값을 조정한다."
))

# Strengths 자리(top-left) -> TN
clean_set_text(find_shape(s7, "TextBox 32"), "정상 예측")
clean_set_text(find_shape(s7, "TextBox 33"), "TRUE NEGATIVE")
clean_set_text(find_shape(s7, "TextBox 34"), "문제 없음 — 정상 거래를 정상으로 정확히 판별")
# Opportunities 자리(bottom-left) -> FN
clean_set_text(find_shape(s7, "TextBox 24"), "탐지 실패")
clean_set_text(find_shape(s7, "TextBox 25"), "FALSE NEGATIVE")
clean_set_text(find_shape(s7, "TextBox 26"), "사기를 정상으로 오판 — 그대로 금전 손실")
# Weaknesses 자리(top-right) -> FP
clean_set_text(find_shape(s7, "TextBox 11"), "오탐지")
clean_set_text(find_shape(s7, "TextBox 12"), "FALSE POSITIVE")
clean_set_text(find_shape(s7, "TextBox 13"), "정상을 사기로 오판 — 고객 재확인 비용(상대적으로 작음)")
# Threats 자리(bottom-right) -> TP
clean_set_text(find_shape(s7, "TextBox 28"), "탐지 성공")
clean_set_text(find_shape(s7, "TextBox 29"), "TRUE POSITIVE")
clean_set_text(find_shape(s7, "TextBox 30"), "사기를 사기로 정확히 판별 — Recall의 분자")

print("S7 built OK")

# ---------------------------------------------------------------- S8 데이터 전처리 & 클리닝
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8, LIGHT)
add_textbox(s8, "3.3 데이터 전처리 & 클리닝", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tbl8 = [
    ["처리", "내용", "근거"],
    ["Time 제거", "거래 순번성 정보, 사기 여부와 직접 관련 없음", "노이즈 제거"],
    ["Amount 로그 변환", "log1p(Amount) → Amount_Scaled 컬럼 추가", "왜도 16.98 → 0.16로 완화"],
    ["V14 이상치 제거", "사기(Class=1) 데이터의 V14 IQR×1.5 범위 밖 값 제거", "상관관계 상위권 변수의 잡음 축소"],
    ["train/test 분할", "7:3, stratify=y", "0.173% 비율을 양쪽 세트에 동일하게 유지"],
]
tbl_w8 = 10200000
add_table_styled(s8, tbl8, CONTENT_LEFT, 1900000, tbl_w8, 4600000,
                  col_widths=[Emu(2600000), Emu(4600000), Emu(3000000)],
                  header_size=14, body_size=13, align_center_cols={0})
add_picture_fit(s8, IMG_DIR + "p23_0.png", CONTENT_LEFT + tbl_w8 + 400000, 1900000,
                CONTENT_W - tbl_w8 - 400000, 4200000)
add_textbox(s8, "사기 클래스는 V14 중앙값이 -6.8로 정상(0 부근)보다 뚜렷이 낮게 분포",
            CONTENT_LEFT + tbl_w8 + 400000, 6200000, CONTENT_W - tbl_w8 - 400000, 700000,
            size=14, color=DARK, alpha=85, align=PP_ALIGN.CENTER, line_spacing=1.2)

print("S8 built OK")

# ---------------------------------------------------------------- S9 샘플링 전략 3종
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9, LIGHT)
add_textbox(s9, "3.5 왜 리샘플링이 필요한가 — 샘플링 전략 3종", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=26, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tbl9 = [
    ["기법", "방식", "학습 데이터 변화"],
    ["언더샘플링", "정상 거래를 사기 건수만큼 무작위로 축소", "199,362건 → 684건 (342/342)"],
    ["SMOTE (채택)", "사기 샘플 주변에 합성 데이터 생성해 오버샘플링", "199,362건 → 398,040건 (199,020/199,020)"],
    ["SMOTE-ENN", "SMOTE 이후 경계가 모호한 샘플을 ENN으로 추가 제거", "199,362건 → 397,717건 (199,020/198,697)"],
]
add_table_styled(s9, tbl9, CONTENT_LEFT, 1750000, CONTENT_W, 2600000,
                  col_widths=[Emu(2400000), Emu(7200000), Emu(6630600)],
                  header_size=14, body_size=13, align_center_cols={0}, highlight_rows={2})
add_picture_fit(s9, IMG_DIR + "p24_0.png", CONTENT_LEFT + 1500000, 4700000, 13230600, 3300000)
add_textbox(s9, "SMOTE는 정상 거래 정보를 버리지 않으면서 사기 샘플을 원본의 약 582배로 증강",
            CONTENT_LEFT, 8200000, CONTENT_W, 500000, size=15, color=DARK, alpha=85,
            align=PP_ALIGN.CENTER)
add_textbox(s9, "채택: SMOTE — 언더샘플링은 정상 정보 손실 과다, SMOTE-ENN은 SMOTE 대비 이점 미미 (§3.6.4)",
            CONTENT_LEFT, 8650000, CONTENT_W, 500000, size=15, bold=True, color=BLUE,
            align=PP_ALIGN.CENTER)

print("S9 built OK")

# ---------------------------------------------------------------- S10 GridSearchCV 튜닝
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10, LIGHT)
add_textbox(s10, "3.5 GridSearchCV 튜닝 — RandomForest", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s10, "Best Params: n_estimators=1000, max_depth=None, min_samples_leaf=1  →  CV ROC-AUC 0.9811",
            CONTENT_LEFT, 1600000, CONTENT_W, 500000, size=17, bold=True, color=BLUE,
            align=PP_ALIGN.CENTER)
tbl10 = [
    ["n_estimators", "max_depth", "min_samples_leaf", "CV ROC-AUC", "순위"],
    ["1000", "None", "1", "0.9811", "1"],
    ["1000", "16", "3", "0.9810", "2"],
    ["1000", "None", "3", "0.9806", "3"],
    ["1000", "16", "1", "0.9806", "4"],
    ["300", "16", "3", "0.9802", "5"],
    ["300", "16", "1", "0.9792", "6"],
    ["300", "None", "3", "0.9791", "7"],
    ["300", "None", "1", "0.9761", "8"],
]
add_table_styled(s10, tbl10, CONTENT_LEFT, 2300000, CONTENT_W, 4600000,
                  header_size=14, body_size=13, highlight_rows={1})
add_textbox(s10, "8개 조합의 CV ROC-AUC가 0.9761~0.9811에 밀집 — 구조적 파라미터보다 임곗값 조정에 튜닝 자원 집중 (§3.6.2)",
            CONTENT_LEFT, 7150000, CONTENT_W, 600000, size=15, color=DARK, alpha=85,
            align=PP_ALIGN.CENTER)

print("S10 built OK")

# ---------------------------------------------------------------- S11 파이프라인 요약 (5단계 템플릿 -> 6단계)
s11 = duplicate_slide(prs, 8)
clean_set_text(find_shape(s11, "TextBox 2"), "3.6 파이프라인 요약")
clean_set_text(find_shape(s11, "TextBox 6"),
                "Time 제거·Amount 로그변환부터 임곗값 조정까지, RandomForest + SMOTE 채택 파이프라인의 전체 흐름이다.")

# template shapes from card 1 (before removal)
tmpl_group = copy.deepcopy(find_shape(s11, "Group 7")._element)
tmpl_icon = copy.deepcopy(find_shape(s11, "Freeform 10")._element)
tmpl_body = copy.deepcopy(find_shape(s11, "TextBox 11")._element)
tmpl_title = copy.deepcopy(find_shape(s11, "TextBox 12")._element)

remove_names = []
for a, b, c, d in [(7, 10, 11, 12), (13, 16, 17, 18), (19, 22, 23, 24), (25, 28, 29, 30), (31, 34, 35, 36)]:
    remove_names += [f"Group {a}", f"Freeform {b}", f"TextBox {c}", f"TextBox {d}"]
for name in remove_names:
    shp = find_shape(s11, name)
    shp._element.getparent().remove(shp._element)

CARD_W6 = 2144954
GAP6 = 323202
START_X = 2076012
steps = [
    ("전처리", "Time 제거\nAmount log1p"),
    ("데이터 분할", "train/test 7:3\n(stratify)"),
    ("이상치 제거", "V14 IQR×1.5\n개별 제거"),
    ("SMOTE", "train에만\n오버샘플링"),
    ("모델 학습", "RandomForest\n(n=1000)"),
    ("임곗값 조정", "재현율 우선\n(0.305)"),
]
spTree = s11.shapes._spTree
for i, (title_txt, body_txt) in enumerate(steps):
    card_left = START_X + i * (CARD_W6 + GAP6)
    center = card_left + CARD_W6 // 2

    # Only the group's OUTER off/ext are resized; the inner freeform/textbox stay in the
    # group's original local (chOff/chExt) coordinate space and are auto-scaled by the
    # renderer to fill the new outer box, so they must NOT be touched here.
    grp = copy.deepcopy(tmpl_group)
    xfrm = grp.find(qn('p:grpSpPr')).find(qn('a:xfrm'))
    off = xfrm.find(qn('a:off'))
    off.set('x', str(card_left))
    ext = xfrm.find(qn('a:ext'))
    ext.set('cx', str(CARD_W6))
    spTree.append(grp)

    icon = copy.deepcopy(tmpl_icon)
    ioff = icon.find(qn('p:spPr')).find(qn('a:xfrm')).find(qn('a:off'))
    ioff.set('x', str(center - 223837))
    spTree.append(icon)

    step_no = add_textbox(s11, f"STEP {i+1}", card_left, 4620000, CARD_W6, 350000,
                           size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    title_sp = copy.deepcopy(tmpl_title)
    toff = title_sp.find(qn('p:spPr')).find(qn('a:xfrm')).find(qn('a:off'))
    toff.set('x', str(center - 929936))
    run = title_sp.find('.//' + qn('a:r'))
    run.find(qn('a:t')).text = title_txt
    spTree.append(title_sp)

    body_sp = copy.deepcopy(tmpl_body)
    boff = body_sp.find(qn('p:spPr')).find(qn('a:xfrm')).find(qn('a:off'))
    boff.set('x', str(center - 870627))
    txBody = body_sp.find(qn('p:txBody'))
    paras = txBody.findall(qn('a:p'))
    lines = body_txt.split("\n")
    base_para = paras[0]
    for extra in paras[1:]:
        txBody.remove(extra)
    base_runs = base_para.findall(qn('a:r'))
    for extra_r in base_runs[1:]:
        base_para.remove(extra_r)
    for li, line in enumerate(lines):
        if li == 0:
            p_el = base_para
        else:
            p_el = copy.deepcopy(base_para)
            txBody.append(p_el)
        r_el = p_el.find(qn('a:r'))
        r_el.find(qn('a:t')).text = line
    spTree.append(body_sp)

print("S11 built OK")

# ---------------------------------------------------------------- S12 성능 지표 결과
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s12, LIGHT)
add_textbox(s12, "3.6 성능 지표 결과 — RandomForest + SMOTE", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=26, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tbl12 = [
    ["임곗값", "Accuracy", "Precision", "Recall", "F1", "ROC-AUC", "AUPRC"],
    ["0.500 (기본값)", "0.9994", "0.8712", "0.7823", "0.8244", "0.9854", "0.8314"],
    ["0.727 (최대 F1)", "-", "0.9558", "0.7347", "0.8308", "0.9854", "0.8314"],
    ["0.305 (채택)", "0.9992", "0.7500", "0.8367", "0.7910", "0.9854", "0.8314"],
]
tbl_w12 = 11400000
add_table_styled(s12, tbl12, CONTENT_LEFT, 1750000, tbl_w12, 2300000,
                  header_size=13, body_size=13, highlight_rows={3})
add_picture_fit(s12, IMG_DIR + "p27_0.png", CONTENT_LEFT + tbl_w12 + 400000, 1750000,
                CONTENT_W - tbl_w12 - 400000, 4200000)
add_textbox(s12, "채택 임곗값(0.305): TN=85,254 · FP=41 · FN=24 · TP=123 — 테스트셋 사기 147건 중 123건(84%) 탐지",
            CONTENT_LEFT, 4450000, tbl_w12, 900000, size=15, color=DARK, alpha=85, line_spacing=1.3)
add_textbox(
    s12,
    "최대 F1(0.727)은 재현율 0.7347에 그쳐 39건을 놓친다. 임곗값을 0.305로 낮추면 F1은 5.9%p 낮아지지만 "
    "재현율은 0.8367로 올라 놓치는 사기가 24건까지 줄어든다 — 정밀도는 0.9558→0.7500으로 하락하지만, "
    "오탐지(FP) 비용이 탐지 실패(FN) 비용보다 훨씬 작다는 원칙에 부합한다.",
    CONTENT_LEFT, 6100000, CONTENT_W, 1600000, size=15, color=DARK, alpha=85, line_spacing=1.35,
)

print("S12 built OK")

# ---------------------------------------------------------------- S13 모델별 비교
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s13, LIGHT)
add_textbox(s13, "3.6 모델별 비교 (SMOTE 적용)", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=27, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tbl13 = [
    ["모델", "Accuracy", "Precision", "Recall", "F1", "ROC-AUC", "AUPRC", "비고"],
    ["LightGBM", "0.9994", "0.8731", "0.7959", "0.8327", "0.9720", "0.8342", "GridSearchCV 튜닝"],
    ["XGBoost", "0.9993", "0.8276", "0.8163", "0.8219", "0.9771", "0.8350", "기본 파라미터"],
    ["RandomForest (채택)", "0.9992", "0.7500", "0.8367", "0.7910", "0.9854", "0.8314", "임곗값 0.305 조정"],
    ["GradientBoost", "0.9993", "0.8438", "0.7347", "0.7855", "0.9749", "0.6687", "임곗값 0.9844 조정"],
    ["LogisticRegression", "0.9994", "0.8550", "0.7619", "0.8057", "0.9675", "0.7024", "임곗값 최적화"],
]
add_table_styled(s13, tbl13, CONTENT_LEFT, 1750000, CONTENT_W, 3400000,
                  header_size=13, body_size=12, highlight_rows={3},
                  col_widths=[Emu(3200000), Emu(1730000), Emu(1730000), Emu(1730000),
                              Emu(1400000), Emu(1500000), Emu(1400000), Emu(3540600)])
add_textbox(
    s13,
    "RandomForest 채택 근거: ROC-AUC 최고(0.9854) + 재현율 최고(0.8367) — AUPRC는 XGBoost·LightGBM이 "
    "근소하게 높지만(0.003~0.004) 이 프로젝트의 '재현율 우선' 목표에는 RandomForest가 최적이다.",
    CONTENT_LEFT, 5450000, CONTENT_W, 1300000, size=16, bold=True, color=BLUE, line_spacing=1.35,
    align=PP_ALIGN.CENTER,
)

print("S13 built OK")

# ---------------------------------------------------------------- S14 3장 종합 결론
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s14, LIGHT)
add_textbox(s14, "3장 종합 결론", CONTENT_LEFT, 700000, CONTENT_W, 750000,
            size=30, bold=True, color=DARK, align=PP_ALIGN.CENTER)
conclusions = [
    ("①", "언더샘플링·SMOTE·SMOTE-ENN 중 정상 거래 정보를 보존하는 SMOTE가 정밀도-재현율 균형이 가장 좋았다"),
    ("②", "SMOTE 적용 후 5개 모델을 비교한 결과 RandomForest가 ROC-AUC(0.9854)와 임곗값 조정 여지에서 가장 우수했다"),
    ("③", "최대 F1 임곗값(0.727) 대신 재현율 우선 임곗값(0.305)을 채택해 테스트셋 사기 147건 중 123건(84%)을 "
          "탐지하고 손실을 24건까지 줄였다"),
]
c_top = 2200000
for i, (num, text) in enumerate(conclusions):
    row_top = c_top + i * 1550000
    add_textbox(s14, num, CONTENT_LEFT, row_top, 900000, 1200000, size=44, bold=True, color=BLUE,
                align=PP_ALIGN.CENTER)
    add_textbox(s14, text, CONTENT_LEFT + 1050000, row_top + 100000, CONTENT_W - 1050000, 1200000,
                size=18, color=DARK, alpha=90, line_spacing=1.3)
add_textbox(s14, "오탐지 비용 < 탐지 실패 비용 — 하나의 원칙을 일관되게 적용한 결과",
            CONTENT_LEFT, c_top + 3 * 1550000 + 150000, CONTENT_W, 600000,
            size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

print("S14 built OK")

# ---------------------------------------------------------------- S15 감사합니다
s15 = duplicate_slide(prs, 9)
clean_set_text(find_shape(s15, "TextBox 2"), "MACHINE LEARNING CASE STUDY")
clean_set_text(find_shape(s15, "TextBox 3"),
                "서울특별시 · ML 종합보고서 4차 · 3장 Credit Card Fraud Detection · 담당 염윤호")
title5b = find_shape(s15, "TextBox 5")
clean_set_text(title5b, "감사합니다")
sub6b = find_shape(s15, "TextBox 6")
clean_set_text(sub6b, "CREDIT CARD FRAUD DETECTION")
sub6b.text_frame.paragraphs[0].runs[0].font.size = Pt(30)

print("S15 built OK")

# =====================================================================
# reorder: remove original 10 template slides, keep only the 15 new ones
# =====================================================================
for _ in range(10):
    remove_slide(prs, 0)

prs.save(OUT_PATH)
print(f"SAVED: {OUT_PATH}  (slide count = {len(prs.slides._sldIdLst)})")





